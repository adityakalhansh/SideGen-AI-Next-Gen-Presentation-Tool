import os
import requests
import urllib.parse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def download_image(prompt, filename):
    try:
        # Pollinations.ai generates images from prompts for free
        encoded_prompt = urllib.parse.quote(prompt)
        url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=800&height=600&nologo=true"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            with open(filename, 'wb') as f:
                f.write(response.content)
            return True
    except Exception as e:
        print(f"Error downloading image: {e}")
    return False

def apply_theme(slide, shape, title_shape, theme_name):
    themes = {
        "dark_modern": {"bg": RGBColor(20, 20, 20), "title": RGBColor(0, 255, 204), "text": RGBColor(240, 240, 240)},
        "minimal_white": {"bg": RGBColor(255, 255, 255), "title": RGBColor(30, 30, 30), "text": RGBColor(60, 60, 60)},
        "gradient_tech": {"bg": RGBColor(40, 10, 60), "title": RGBColor(255, 100, 255), "text": RGBColor(230, 230, 250)},
        "startup_pitch": {"bg": RGBColor(250, 245, 235), "title": RGBColor(255, 80, 80), "text": RGBColor(40, 40, 40)},
        "education_classic": {"bg": RGBColor(240, 248, 255), "title": RGBColor(0, 102, 204), "text": RGBColor(50, 50, 50)}
    }
    
    theme = themes.get(theme_name, themes["minimal_white"])
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["bg"]
    
    # Set title color
    if title_shape and title_shape.has_text_frame:
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = theme["title"]
                run.font.name = "Arial" # Standard safe font
                
    # Set text color
    if shape and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = theme["text"]
                run.font.name = "Arial"

def create_presentation(slides_data, theme, output_filename):
    prs = Presentation()
    
    # Create temp directory for images
    if not os.path.exists("temp_images"):
        os.makedirs("temp_images")
        
    for i, slide_info in enumerate(slides_data):
        title = slide_info.get("title", "")
        bullets = slide_info.get("bullets", [])
        img_prompt = slide_info.get("image_prompt", f"presentation slide about {title}")
        
        # Decide layout
        if i == 0:
            # Title slide
            slide_layout = prs.slide_layouts[0] 
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            if bullets:
                subtitle_shape.text = bullets[0]
            
            apply_theme(slide, subtitle_shape, title_shape, theme)
            
        else:
            # Content slide with image
            # We use a Title and Content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
            
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            
            # Download and add image
            img_filename = f"temp_images/slide_{i}.jpg"
            has_image = download_image(img_prompt, img_filename)
            
            if has_image:
                # Reduce width of the text box to make room for the image
                body_shape.width = Inches(4.5)
                img_left = Inches(5.2)
                img_top = Inches(2.0)
                img_width = Inches(4.3)
                try:
                    slide.shapes.add_picture(img_filename, img_left, img_top, width=img_width)
                except Exception as e:
                    print(f"Failed to add picture {img_filename}: {e}")
                    # If picture fails, restore text box width
                    body_shape.width = Inches(9.0)
            
            tf.word_wrap = True
            for j, bullet in enumerate(bullets):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                
            apply_theme(slide, body_shape, title_shape, theme)

    prs.save(output_filename)
    
    # Cleanup temp images
    for file in os.listdir("temp_images"):
        os.remove(os.path.join("temp_images", file))
    try:
        os.rmdir("temp_images")
    except:
        pass
    
    return output_filename
